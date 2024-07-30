import os
import shutil
import threading
import tkinter as tk
from tkinter import filedialog, messagebox
from PyPDF2 import PdfReader
import pandas as pd
from docx import Document
from pptx import Presentation
import time

# Variable global para controlar el proceso de búsqueda
stop_search = False

def search_keywords_in_file(file_path, keywords):
    ext = os.path.splitext(file_path)[1].lower()
    content = ""

    try:
        if ext == ".pdf":
            with open(file_path, "rb") as f:
                reader = PdfReader(f)
                for page in reader.pages:
                    content += page.extract_text()
        elif ext in [".xls", ".xlsx", ".xlsm", ".xltm"]:
            df = pd.read_excel(file_path, engine='openpyxl')
            content = df.to_string()
        elif ext == ".csv":
            df = pd.read_csv(file_path, encoding='latin1', low_memory=False)  # Cambia la codificación si es necesario
            content = df.to_string()
        elif ext == ".docx":
            doc = Document(file_path)
            for para in doc.paragraphs:
                content += para.text
        elif ext == ".txt":
            with open(file_path, "r", encoding="utf-8") as f:
                content = f.read()
        elif ext == ".pptx":
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content += shape.text
    except Exception as e:
        print(f"Error reading {file_path}: {e}")
        return False

    # Normalizar el contenido para la búsqueda
    content = content.lower()

    # Incluir el nombre del archivo en la búsqueda
    file_name = os.path.basename(file_path).lower()

    for keyword in keywords:
        if keyword.lower() in content or keyword.lower() in file_name:
            return True
    return False

def copy_files_with_keywords(src_folder, dest_folder, keywords):
    global stop_search
    copied_files = []
    file_count = 0
    start_time = time.time()
    for root_dir, _, files in os.walk(src_folder):
        if stop_search:
            break
        for file in files:
            if stop_search:
                break
            file_path = os.path.join(root_dir, file)
            if search_keywords_in_file(file_path, keywords):
                dest_path = os.path.join(dest_folder, file)
                shutil.copy(file_path, dest_path)
                copied_files.append(file)
                file_count += 1
                result_list.insert(tk.END, file)
                result_list.yview(tk.END)
                file_count_label.config(text=f"Archivos Encontrados: {file_count}")
                root.update_idletasks()  # Actualiza la GUI
    end_time = time.time()
    elapsed_time = end_time - start_time
    time_label.config(text=f"Tiempo Transcurrido: {elapsed_time:.2f} segundos")
    return copied_files

def browse_folder(entry):
    folder_selected = filedialog.askdirectory()
    entry.delete(0, tk.END)
    entry.insert(0, folder_selected)

def start_search():
    global stop_search
    stop_search = False
    src_folder = src_entry.get()
    dest_folder = dest_entry.get()
    keywords = keyword_entry.get().split(",")

    if not src_folder or not dest_folder or not keywords:
        messagebox.showwarning("Error de Entrada", "Por favor, complete todos los campos.")
        return

    result_list.delete(0, tk.END)
    file_count_label.config(text="Archivos Encontrados: 0")
    time_label.config(text="Tiempo Transcurrido: 0.00 segundos")
    search_thread = threading.Thread(target=copy_files_with_keywords, args=(src_folder, dest_folder, keywords))
    search_thread.start()

def stop_search_process():
    global stop_search
    stop_search = True

# Configuración de la GUI
root = tk.Tk()
root.title("Buscador de Palabras Clave en Archivos")

tk.Label(root, text="Carpeta de Origen:").grid(row=0, column=0, padx=10, pady=5)
src_entry = tk.Entry(root, width=50)
src_entry.grid(row=0, column=1, padx=10, pady=5)
tk.Button(root, text="Buscar", command=lambda: browse_folder(src_entry)).grid(row=0, column=2, padx=10, pady=5)

tk.Label(root, text="Carpeta de Destino:").grid(row=1, column=0, padx=10, pady=5)
dest_entry = tk.Entry(root, width=50)
dest_entry.grid(row=1, column=1, padx=10, pady=5)
tk.Button(root, text="Buscar", command=lambda: browse_folder(dest_entry)).grid(row=1, column=2, padx=10, pady=5)

tk.Label(root, text="Palabras Clave (separadas por comas):").grid(row=2, column=0, padx=10, pady=5)
keyword_entry = tk.Entry(root, width=50)
keyword_entry.grid(row=2, column=1, padx=10, pady=5)

tk.Button(root, text="Iniciar Búsqueda", command=start_search).grid(row=3, column=0, pady=10)
tk.Button(root, text="Detener Búsqueda", command=stop_search_process).grid(row=3, column=1, pady=10)

file_count_label = tk.Label(root, text="Archivos Encontrados: 0")
file_count_label.grid(row=4, column=0, padx=10, pady=5)

time_label = tk.Label(root, text="Tiempo Transcurrido: 0.00 segundos")
time_label.grid(row=4, column=1, padx=10, pady=5)

tk.Label(root, text="Archivos Copiados:").grid(row=5, column=0, padx=10, pady=5)
result_list = tk.Listbox(root, width=70, height=10)
result_list.grid(row=5, column=1, columnspan=2, padx=10, pady=5)

# Agregar el label con la información de contacto
contact_label = tk.Label(root, text="Desarrollado por Abel Ricardo Ibarra\n+528787015229 - abel.ricardo@data146.com", font=("Arial", 10), fg="blue")
contact_label.grid(row=6, column=0, columnspan=3, pady=20)

root.mainloop()